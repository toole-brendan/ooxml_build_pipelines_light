#!/usr/bin/env python3
"""
build_deck_pptx.py -- Build Saronic / Port Alpha MRO Entry Due Diligence Deck.

Generates a 10-slide PowerPoint scaffold matching the spec in DECK_PROPOSED.md.
Each chart slide carries a data block formatted for think-cell's datasheet
paste workflow; to create the final chart, select the block in PowerPoint and
insert a think-cell chart of the type noted above the block.

Usage:
    python3 deck/build_deck_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ---------- Slide geometry (16:9 widescreen) ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.55)
MARGIN_R = Inches(0.55)
USABLE_W = SLIDE_W - MARGIN_L - MARGIN_R  # 12.233"

TITLE_TOP = Inches(0.4)
TITLE_H = Inches(0.85)
TITLE_RULE_Y = Inches(1.3)

BODY_TOP = Inches(1.45)
BODY_BOTTOM = Inches(6.95)
BODY_H = BODY_BOTTOM - BODY_TOP

FOOTER_TOP = Inches(7.0)
FOOTER_H = Inches(0.45)

# Left-right body split (1.15 : 1 with 0.3" gap, per HTML)
L_PANEL_X = MARGIN_L
L_PANEL_W = Inches(6.5)
R_PANEL_X = MARGIN_L + L_PANEL_W + Inches(0.3)
R_PANEL_W = SLIDE_W - R_PANEL_X - MARGIN_R


# ---------- Palette (from HTML CSS) ----------
NAVY_900 = RGBColor(0x1F, 0x2D, 0x44)
NAVY_800 = RGBColor(0x24, 0x36, 0x50)
NAVY_700 = RGBColor(0x2C, 0x42, 0x62)
NAVY_600 = RGBColor(0x36, 0x52, 0x7A)
SLATE_500 = RGBColor(0x6F, 0x81, 0x97)
SLATE_400 = RGBColor(0x8A, 0x9B, 0xB0)
SLATE_300 = RGBColor(0xB6, 0xC2, 0xD1)
SLATE_200 = RGBColor(0xD6, 0xDD, 0xE6)
SLATE_100 = RGBColor(0xE7, 0xEC, 0xF2)
LIGHTBLUE_100 = RGBColor(0xDB, 0xE6, 0xF1)
LIGHTBLUE_200 = RGBColor(0xC5, 0xD6, 0xE6)
LIGHTBLUE_300 = RGBColor(0xA8, 0xC0, 0xD6)
LIGHTBLUE_400 = RGBColor(0x7F, 0xA0, 0xBF)
INK = RGBColor(0x2A, 0x33, 0x40)
INK_SOFT = RGBColor(0x4B, 0x56, 0x66)
RULE = RGBColor(0xC9, 0xD2, 0xDC)
RULE_SOFT = RGBColor(0xE3, 0xE8, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = WHITE

FONT_SANS = "Helvetica Neue"


# ---------- Shape helpers ----------
def _set_text(tf, text, *, size=10, color=INK, bold=False, italic=False,
              align=PP_ALIGN.LEFT, font=FONT_SANS, anchor=MSO_ANCHOR.TOP):
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color


def _add_run(paragraph, text, *, size=10, color=INK, bold=False, italic=False,
             font=FONT_SANS):
    run = paragraph.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return run


def add_textbox(slide, left, top, width, height, *, fill=None, line=None,
                line_width=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.shadow.inherit = False
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        if line_width is not None:
            box.line.width = line_width
    return box


def add_hline(slide, left, top, width, color=RULE, weight_pt=0.75):
    ln = slide.shapes.add_connector(1, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight_pt)
    return ln


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------- Title / Subtitle / Footer ----------
def add_title(slide, topic, assertion):
    tb = slide.shapes.add_textbox(MARGIN_L, TITLE_TOP, USABLE_W, TITLE_H)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, topic, size=18, color=NAVY_900)
    _add_run(p, "   |   ", size=18, color=SLATE_400)
    _add_run(p, assertion, size=18, color=INK)
    add_hline(slide, MARGIN_L, TITLE_RULE_Y, USABLE_W, color=RULE, weight_pt=0.75)


def add_subtitle(slide, left, top, width, text):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    _set_text(tb.text_frame, text, size=11, color=INK_SOFT, italic=True)
    return tb


def add_footer(slide, text):
    tb = slide.shapes.add_textbox(MARGIN_L, FOOTER_TOP, USABLE_W, FOOTER_H)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_run(p, text, size=7.5, color=SLATE_500)
    add_hline(slide, MARGIN_L, FOOTER_TOP - Inches(0.05), USABLE_W,
              color=RULE_SOFT, weight_pt=0.5)


# ---------- Callout (light-blue italic) ----------
def add_callout(slide, left, top, width, height, lead, body):
    box = add_textbox(slide, left, top, width, height,
                      fill=LIGHTBLUE_100, line=LIGHTBLUE_200, line_width=Pt(0.5))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if lead:
        _add_run(p, lead + " ", size=8.5, color=NAVY_900, bold=True)
    _add_run(p, body, size=8.5, color=NAVY_900, italic=True)
    return box


# ---------- Data table ----------
def _cell_border(cell, side, color=RULE_SOFT, weight_pt=0.5):
    tag = {"top": "a:lnT", "bottom": "a:lnB", "left": "a:lnL", "right": "a:lnR"}[side]
    tcPr = cell._tc.get_or_add_tcPr()
    for old in tcPr.findall(qn(tag)):
        tcPr.remove(old)
    ln = etree.SubElement(tcPr, qn(tag))
    ln.set("w", str(int(Pt(weight_pt))))
    ln.set("cap", "flat")
    ln.set("cmpd", "sng")
    ln.set("algn", "ctr")
    fill = etree.SubElement(ln, qn("a:solidFill"))
    srgb = etree.SubElement(fill, qn("a:srgbClr"))
    srgb.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))
    etree.SubElement(ln, qn("a:prstDash")).set("val", "solid")
    etree.SubElement(ln, qn("a:round"))


def _cell_no_border(cell, side):
    tag = {"top": "a:lnT", "bottom": "a:lnB", "left": "a:lnL", "right": "a:lnR"}[side]
    tcPr = cell._tc.get_or_add_tcPr()
    for old in tcPr.findall(qn(tag)):
        tcPr.remove(old)
    ln = etree.SubElement(tcPr, qn(tag))
    etree.SubElement(ln, qn("a:noFill"))


def _fill_cell(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _write_cell(cell, text, *, size=9, color=INK, bold=False, italic=False,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    _set_text(cell.text_frame, text, size=size, color=color, bold=bold,
              italic=italic, align=align, anchor=anchor)
    cell.margin_left = Inches(0.07)
    cell.margin_right = Inches(0.07)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)


def add_data_table(slide, left, top, width, height, headers, rows, *,
                   col_widths=None, right_align_cols=None,
                   header_size=9, body_size=8.5, row_height=Inches(0.24)):
    """
    rows: list of dicts with keys:
        'cells': [str, ...]
        'kind': 'normal' | 'subhead' | 'total' | 'sub'
        'bolds': optional [bool] per cell
    """
    right_align_cols = set(right_align_cols or [])
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(width * w / total))

    tbl.rows[0].height = Inches(0.3)
    for i, h in enumerate(headers):
        c = tbl.cell(0, i)
        _fill_cell(c, WHITE)
        align = PP_ALIGN.RIGHT if i in right_align_cols else PP_ALIGN.LEFT
        _write_cell(c, h, size=header_size, color=NAVY_900, bold=True, align=align)
        for side in ("top", "left", "right"):
            _cell_no_border(c, side)
        _cell_border(c, "bottom", color=NAVY_800, weight_pt=1.0)

    for r_idx, row in enumerate(rows, start=1):
        kind = row.get("kind", "normal")
        tbl.rows[r_idx].height = row_height
        if kind == "subhead":
            c = tbl.cell(r_idx, 0)
            if n_cols > 1:
                c.merge(tbl.cell(r_idx, n_cols - 1))
            _fill_cell(c, LIGHTBLUE_100)
            _write_cell(c, row["cells"][0], size=body_size - 1.5, color=NAVY_900,
                        bold=True, align=PP_ALIGN.LEFT)
            for side in ("top", "left", "right"):
                _cell_no_border(c, side)
            _cell_border(c, "bottom", color=LIGHTBLUE_200, weight_pt=0.5)
            continue
        for j, val in enumerate(row["cells"]):
            c = tbl.cell(r_idx, j)
            _fill_cell(c, WHITE)
            align = PP_ALIGN.RIGHT if j in right_align_cols else PP_ALIGN.LEFT
            bold = False
            italic = False
            color = INK
            size = body_size
            if kind == "total":
                bold = True
                color = NAVY_900
            elif kind == "sub":
                italic = True
                color = INK_SOFT
            bolds = row.get("bolds")
            if bolds and j < len(bolds) and bolds[j]:
                bold = True
                color = NAVY_900
            _write_cell(c, val, size=size, color=color, bold=bold, italic=italic,
                        align=align)
            for side in ("top", "left", "right"):
                _cell_no_border(c, side)
            if kind == "total":
                _cell_border(c, "top", color=NAVY_800, weight_pt=0.75)
                _cell_no_border(c, "bottom")
            else:
                _cell_border(c, "bottom", color=RULE_SOFT, weight_pt=0.4)
    return tbl_shape


# ---------- Think-cell data block ----------
def add_think_cell_block(slide, left, top, width, height, chart_type, headers,
                         rows, *, col_widths=None, right_align_cols=None,
                         note=None):
    """Render a think-cell-ready data block with an instruction header."""
    header_h = Inches(0.55)
    instr = add_textbox(slide, left, top, width, header_h,
                        fill=LIGHTBLUE_100, line=LIGHTBLUE_200, line_width=Pt(0.5))
    tf = instr.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    _add_run(p, "Think-cell ", size=9, color=NAVY_900, bold=True)
    _add_run(p, chart_type, size=9, color=NAVY_900, bold=True)
    p2 = tf.add_paragraph()
    _add_run(p2, "Select this data block, then Insert -> think-cell -> "
             + chart_type + " and paste into the think-cell datasheet.",
             size=8, color=INK_SOFT, italic=True)
    if note:
        p3 = tf.add_paragraph()
        _add_run(p3, note, size=7.5, color=INK_SOFT, italic=True)

    tbl_top = top + header_h + Inches(0.05)
    tbl_h = height - header_h - Inches(0.05)
    row_h = Inches(0.22) if len(rows) > 8 else Inches(0.26)
    add_data_table(slide, left, tbl_top, width, tbl_h, headers, rows,
                   col_widths=col_widths, right_align_cols=right_align_cols,
                   header_size=8.5, body_size=8, row_height=row_h)


# ---------- Slide 1: Research Context & Approach ----------
def slide_01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s,
              "Research Context & Approach",
              "Sizing U.S. Navy and Coast Guard MRO contracting activity to "
              "support the new shipyard's business-line decision")

    col_w = (USABLE_W - Inches(0.5)) / 2
    col1_x = MARGIN_L
    col2_x = MARGIN_L + col_w + Inches(0.5)

    def _text_col(x, header, bullets):
        hb = slide_add_section_header(s, x, BODY_TOP, col_w, header)
        y = BODY_TOP + Inches(0.55)
        for b in bullets:
            tb = s.shapes.add_textbox(x, y, col_w, Inches(0.6))
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.0)
            tf.margin_bottom = Inches(0.0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            _add_run(p, "- ", size=11, color=SLATE_500)
            _add_run(p, b, size=11, color=INK)
            y += Inches(0.65)

    _text_col(col1_x, "Context", [
        "Saronic is finalizing dimensions for a new U.S. shipyard.",
        "Candidate sites under evaluation: Gulf Coast (Brownsville, TX) and "
        "California.",
        "Leadership is deciding whether to include U.S. Navy and U.S. Coast "
        "Guard MRO contracting among the yard's business lines.",
        "This research quantifies the MRO market and characterizes its "
        "structure so the business-line decision can be informed by FY2025 "
        "market data rather than planning assumptions.",
    ])
    _text_col(col2_x, "Objectives", [
        "Size the FY2025 U.S. Navy and U.S. Coast Guard MRO contracting "
        "market (TAM) using FPDS contract obligations.",
        "Decompose TAM by work segment, vessel category, prime contractor, "
        "Regional Maintenance Center (RMC) geography, and funding "
        "appropriation.",
        "Identify the structural characteristics of the largest segment "
        "(depot ship repair) and its entry structure.",
        "Reconcile the TAM against adjacent ship-dollar pools (newbuild, "
        "reactor procurement, public-yard labor) that are outside MRO scope "
        "but potentially relevant context.",
    ])

    divider_x = MARGIN_L + col_w + Inches(0.25)
    ln = s.shapes.add_connector(1, divider_x, BODY_TOP, divider_x, BODY_BOTTOM)
    ln.line.color.rgb = RULE
    ln.line.width = Pt(0.5)


def slide_add_section_header(slide, x, y, w, text):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.45))
    _set_text(tb.text_frame, text, size=15, color=NAVY_900, bold=True)
    add_hline(slide, x, y + Inches(0.45), w, color=NAVY_800, weight_pt=0.75)
    return tb


# ---------- Slide 2: TAM & Scope ----------
def slide_02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "TAM & Scope",
              "FY2025 Navy and Coast Guard MRO TAM totaled $7.1B")

    # Funnel (left)
    funnel_w = L_PANEL_W * 0.78
    funnel_x = L_PANEL_X + (L_PANEL_W - funnel_w) / 2
    y = BODY_TOP + Inches(0.1)
    box_h = Inches(0.75)
    arrow_h = Inches(0.28)

    def _funnel_box(tier, big, sub, fill, color):
        nonlocal y
        b = add_textbox(s, funnel_x, y, funnel_w, box_h, fill=fill)
        tf = b.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_run(p, big, size=15, color=color, bold=True)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _add_run(p2, sub, size=9.5, color=color, italic=True)
        y += box_h

    def _funnel_arrow(label):
        nonlocal y
        tb = s.shapes.add_textbox(funnel_x, y, funnel_w, arrow_h)
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_run(p, "v  ", size=12, color=SLATE_500, bold=True)
        _add_run(p, label, size=9, color=INK_SOFT, italic=True)
        y += arrow_h

    _funnel_box("t1", "2,539 Active PSCs",
                "Every federal contract-action code category",
                LIGHTBLUE_100, NAVY_900)
    _funnel_arrow("USN / USCG only, Services codes only")
    _funnel_box("t2", "~1,800 Services-class PSCs",
                "Services-only codes (J, K, N, H, L, M families); products excluded",
                LIGHTBLUE_200, NAVY_900)
    _funnel_arrow("MRO codes only")
    _funnel_box("t3", "65 Ship MRO PSCs",
                "Navy + USCG ship repair, modification, overhaul and husbanding codes",
                NAVY_800, WHITE)
    _funnel_arrow("")
    _funnel_box("term", "FY2025 MRO TAM = $7.1B", "",
                NAVY_900, WHITE)

    # Note/Detail table (right)
    add_data_table(s, R_PANEL_X, BODY_TOP + Inches(0.3), R_PANEL_W,
                   Inches(5.0),
                   ["Note", "Detail"],
                   [
                       {"cells": ["PSC",
                        "A 4-character Product and Service Code assigned by "
                        "the contracting officer at the buying command to "
                        "every federal contract action."],
                        "bolds": [True, False]},
                       {"cells": ["Awards data vs. budget materials",
                        "FPDS obligations capture executed contract-level "
                        "spend, attributing dollars to vendor, work type, and "
                        "hull program. Navy and Coast Guard budget exhibits "
                        "report planned maintenance at the program level, "
                        "with no contractor attribution and no PSC-level "
                        "work-type detail."],
                        "bolds": [True, False]},
                       {"cells": ["In scope",
                        "All U.S. Navy and U.S. Coast Guard vessel types. "
                        "Work types: depot ship repair, equipment "
                        "maintenance, modification, installation, QC and "
                        "inspection, OEM technical representation, and "
                        "husbanding services."],
                        "bolds": [True, False]},
                       {"cells": ["Out of scope",
                        "In-house labor at the four public naval shipyards "
                        "(Portsmouth, Norfolk, Puget Sound, and Pearl "
                        "Harbor); RCOH and reactor-plant sustainment "
                        "(bundled under shipbuilding PSCs); newbuild and "
                        "product procurement."],
                        "bolds": [True, False]},
                   ],
                   col_widths=[30, 70], header_size=11, body_size=9,
                   row_height=Inches(1.1))

    add_footer(s, "Source: FPDS FY2025 contract obligations (U.S. Navy + "
               "U.S. Coast Guard, 65 services PSCs, post-exclusions). "
               "Data as of April 2026.")


# ---------- Slide 3: Addressable vs Adjacent Spend ----------
def slide_03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Addressable vs Adjacent Spend",
              "The $7.1B MRO TAM is the addressable slice of a much larger "
              "FY2025 Navy and Coast Guard ship-dollar pool; ~$38B newbuild, "
              "~$9.5B implied public-yard labor, and ~$1.9B reactor product "
              "sit adjacent to but outside the MRO scope")

    # LEFT: waterfall data block
    add_subtitle(s, L_PANEL_X, BODY_TOP, L_PANEL_W,
                 "FY2025 Navy + USCG ship-dollar pools bridged to addressable "
                 "MRO TAM ($M)")
    add_think_cell_block(
        s, L_PANEL_X, BODY_TOP + Inches(0.32), L_PANEL_W, Inches(4.5),
        "Waterfall",
        ["Step", "$M", "Type"],
        [
            {"cells": ["FY25 Navy + USCG ship-dollar pool", "56,577", "e (endpoint)"]},
            {"cells": ["- PSC 1905 Newbuild (all platforms)", "-38,100", "delta"]},
            {"cells": ["- PSC 4470 Nuclear Reactors", "-1,875", "delta"]},
            {"cells": ["- Implied public-yard labor (OMN 1B4B)", "-9,535", "delta"]},
            {"cells": ["- Nuclear MRO PSCs (J044/K044/N044)", "0", "delta"]},
            {"cells": ["FY25 MRO TAM (addressable)", "7,067", "e (endpoint)"],
             "kind": "total"},
        ],
        col_widths=[55, 20, 25], right_align_cols={1},
        note="In the think-cell datasheet, set bars 1 and 6 as endpoints "
             "(prefix value with 'e') and bars 2-5 as deltas.")

    # RIGHT: budget anchors table
    tbl_top = BODY_TOP + Inches(0.05)
    add_data_table(s, R_PANEL_X, tbl_top, R_PANEL_W, Inches(3.7),
                   ["Line Item", "FY25 $M", "Source"],
                   [
                       {"cells": ["OMN - Ship Maintenance (SAG 1B4B)"], "kind": "subhead"},
                       {"cells": ["Total 1B4B Ship Maintenance", "11,764",
                        "OMN PB Vol 1, p. 153"]},
                       {"cells": ["   of which Ship Maintenance By Contract "
                        "(CE 928)", "2,228", "OMN PB Vol 1, p. 178 (OP-32 L928)"],
                        "kind": "sub"},
                       {"cells": ["Implied non-contract (public-yard labor + supt.)",
                        "9,535", "Computed (1B4B minus CE 928)"]},
                       {"cells": ["SCN - Nuclear-Platform Line Items"], "kind": "subhead"},
                       {"cells": ["LI 1045 Columbia Class (construction + CY AP)",
                        "9,581", "SCN P-40, pp. 31 + 57"]},
                       {"cells": ["LI 2013 Virginia Class (construction + CY AP)",
                        "13,221", "SCN P-40, pp. 155 + 171"]},
                       {"cells": ["LI 2086 CVN Refueling Overhauls (FY25 net)",
                        "1,480", "SCN P-40, p. 175"]},
                       {"cells": ["LI 2001 Carrier Replacement (CVN-80)",
                        "1,359", "SCN P-40, p. 73"]},
                       {"cells": ["LI 2004 CVN-81", "675", "SCN P-40, p. 117"]},
                       {"cells": ["SCN nuclear-platform subtotal", "26,316", "Sum"],
                        "kind": "total"},
                       {"cells": ["Nuclear MRO PSC Emptiness (J044 / K044 / N044)"],
                        "kind": "subhead"},
                       {"cells": ["Nuclear Maintenance / Modification / Installation",
                        "~0", "Bundled under shipbuilding PSCs"]},
                   ],
                   col_widths=[52, 16, 32], right_align_cols={1},
                   header_size=8.5, body_size=7.5, row_height=Inches(0.22))

    add_callout(
        s, R_PANEL_X, Inches(5.1), R_PANEL_W, Inches(1.15),
        "Note:",
        "PSC 1905 is not purely newbuild. Individual depot events on "
        "nuclear-platform hulls (e.g., the $424M HII USS Boise SSN-764 "
        "Engineered Overhaul at Newport News) are coded as shipbuilding rather "
        "than J998/J999. The Services MRO TAM on Slide 2 therefore undercounts "
        "private-sector sub and carrier depot work that is bundled under the "
        "newbuild PSC.")

    add_footer(s,
               "Notes: (1) Public-yard labor implied: OMN 1B4B total ($11,764M) "
               "minus Ship Maintenance By Contract CE 928 ($2,228M). (2) The "
               "four public naval shipyards (Portsmouth, Norfolk, Puget Sound, "
               "Pearl Harbor) are NWCF activities -- payroll generates no FPDS "
               "record. (3) J044/K044/N044 ~$0 because reactor work is "
               "contracted under ship-level PSC 1905 and PSC 4470. "
               "Sources: FPDS FY2025 contract obligations; FY2026 President's "
               "Budget exhibits (OMN Vol 1 SAG 1B4B OP-5, SCN P-1). "
               "Data as of April 2026.")


# ---------- Slide 4: Work Segments ----------
def slide_04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Work Segments",
              "Depot ship repair made up ~68% of FY2025 MRO TAM, while "
              "adjacent sustainment segments were smaller and more specialized")

    add_subtitle(s, L_PANEL_X, BODY_TOP, L_PANEL_W,
                 "FY2025 MRO TAM by Work Segment ($M)")
    add_think_cell_block(
        s, L_PANEL_X, BODY_TOP + Inches(0.32), L_PANEL_W, Inches(3.3),
        "100% Stacked Column (single column)",
        ["Segment", "FY25 $M", "% of TAM"],
        [
            {"cells": ["Depot Ship Repair", "4,781", "68%"], "bolds": [True, False, False]},
            {"cells": ["Hull, Mechanical & Electrical (HM&E)", "938", "13%"]},
            {"cells": ["Combat Systems Sustainment", "585", "8%"]},
            {"cells": ["Port & Technical Services", "431", "6%"]},
            {"cells": ["Electronics & C4ISR Sustainment", "333", "5%"]},
            {"cells": ["Total FY25 MRO TAM", "7,068", "100%"], "kind": "total"},
        ],
        col_widths=[55, 22, 23], right_align_cols={1, 2},
        note="One x-axis category (FY25 MRO TAM); stack the five segments in "
             "the order above. Set chart to 100% stacked.")

    add_data_table(s, R_PANEL_X, BODY_TOP, R_PANEL_W, Inches(3.5),
                   ["Work Segment", "Coverage"],
                   [
                       {"cells": ["Depot Ship Repair",
                        "Whole-ship maintenance availabilities at Pacific and "
                        "Atlantic Regional Maintenance Centers, awarded "
                        "through MAC-MO IDIQ task orders."],
                        "bolds": [True, False]},
                       {"cells": ["Hull, Mechanical & Electrical (HM&E)",
                        "Propulsion accessories, pumps, valves, piping, HVAC, "
                        "diesel engines, and ship structural systems."],
                        "bolds": [True, False]},
                       {"cells": ["Combat Systems Sustainment",
                        "Weapons, fire control, VLS, guided missiles "
                        "(including Trident II sustainment on Ohio-class "
                        "SSBNs), and aircraft launch/arresting equipment."],
                        "bolds": [True, False]},
                       {"cells": ["Port & Technical Services",
                        "Quality control, testing and inspection; OEM "
                        "technical representation; husbanding (fuel, "
                        "transport, port visits); and shipyard operations "
                        "support."],
                        "bolds": [True, False]},
                       {"cells": ["Electronics & C4ISR Sustainment",
                        "Afloat C4ISR -- radar, sonar, radio and network "
                        "systems, navigation, alarms, and electrical signal "
                        "equipment."],
                        "bolds": [True, False]},
                   ],
                   col_widths=[32, 68], header_size=9, body_size=8,
                   row_height=Inches(0.65))

    add_callout(
        s, R_PANEL_X, Inches(5.35), R_PANEL_W, Inches(0.95),
        "Note:",
        "The $4.8B depot segment is awarded through a single three-tier "
        "structure -- MSRA pre-qualification -> MAC-MO IDIQ -> fixed-price "
        "task orders against third-party planner specs. Entry requires "
        "qualification at all three levels. Slide 5 decomposes the depot "
        "segment by scope group and contractor tier.")

    add_footer(s,
               "Notes: (1) Segment detail sums to $7,068M; TAM headline "
               "reports $7,067M -- figures may not sum to headline due to "
               "rounding. (2) Nuclear Propulsion Sustainment PSCs "
               "(J044/K044/N044) appear at ~$0 in FY25 because reactor work "
               "is contracted under ship-level shipbuilding codes at HII "
               "Newport News, Fluor Marine Propulsion, and Bechtel -- not "
               "standalone services PSCs. "
               "Source: FPDS FY2025 contract obligations (U.S. Navy + U.S. "
               "Coast Guard, 65 services PSCs, post-exclusions). Data as of "
               "April 2026.")


# ---------- Slide 5: Depot Ship Repair Deep Dive ----------
def slide_05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Depot Ship Repair Deep Dive",
              "Six CONUS Tier-1 yards captured ~60% of the $4.9B FY2025 depot "
              "segment; 65% of depot $ is full-ship availabilities awarded "
              "through a three-tier MSRA / MAC-MO prime structure")

    add_subtitle(s, L_PANEL_X, BODY_TOP, L_PANEL_W,
                 "FY2025 J998 / J999 depot obligations - IDV scope group x "
                 "contractor tier ($M)")
    # Mekko: first row = column widths; first col = row labels; body = $
    add_think_cell_block(
        s, L_PANEL_X, BODY_TOP + Inches(0.32), L_PANEL_W, Inches(4.35),
        "Mekko (Marimekko)",
        ["Tier", "Full-Ship", "MSC Avail", "FDNF MSRA", "Other/Supt",
         "USCG Cut.", "Trade IDIQ", "P&E"],
        [
            {"cells": ["(column width $M)", "3,258", "551", "513", "306",
                      "144", "123", "108"], "kind": "sub"},
            {"cells": ["Tier 1 CONUS Complex", "2,867", "386", "0", "92",
                      "72", "31", "11"]},
            {"cells": ["Tier 2 Regional", "195", "99", "0", "86",
                      "50", "68", "19"]},
            {"cells": ["Tier 3 Technical Services", "98", "44", "0", "107",
                      "22", "25", "78"]},
            {"cells": ["Tier 4 FDNF Foreign Yard", "98", "22", "513", "0",
                      "0", "0", "0"]},
            {"cells": ["Other", "0", "0", "0", "21", "0", "0", "0"]},
            {"cells": ["Column total", "3,258", "551", "513", "306",
                      "144", "123", "108"], "kind": "total"},
        ],
        col_widths=[24, 11, 11, 11, 11, 11, 11, 10],
        right_align_cols={1, 2, 3, 4, 5, 6, 7},
        note="In the think-cell datasheet: first row becomes the column-width "
             "series; each Tier row is a stack series. Set data type to "
             "absolute values; Mekko normalizes stacks to 100% per column.")

    add_data_table(s, R_PANEL_X, BODY_TOP, R_PANEL_W, Inches(3.05),
                   ["Tier 1 CONUS Prime", "Primary Yards", "Primary RMC",
                    "FY25 $M"],
                   [
                       {"cells": ["BAE Ship Repair",
                        "San Diego, Norfolk, Jacksonville",
                        "SWRMC + MARMC + SERMC", "~1,073"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["General Dynamics (NASSCO + Continental "
                        "Maritime)", "San Diego", "SWRMC", "~640"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["HII (Continental Maritime + Metro Machine "
                        "+ MHI)", "San Diego, Norfolk", "SWRMC + MARMC", "~390"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Vigor", "Portland OR, Seattle WA", "NW RMC",
                        "~440"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Detyens", "Charleston", "SERMC", "~225"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["NASSCO Mayport + others", "Mayport, Kings Bay",
                        "SERMC", "~100"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Tier 1 CONUS subtotal", "", "", "~2,868"],
                        "kind": "total"},
                   ],
                   col_widths=[38, 32, 18, 12], right_align_cols={3},
                   header_size=8.5, body_size=7.5, row_height=Inches(0.34))

    # Entry structure box
    entry = add_textbox(s, R_PANEL_X, Inches(4.8), R_PANEL_W, Inches(1.15),
                        fill=SLATE_100, line=RULE, line_width=Pt(0.5))
    tf = entry.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    _add_run(p, "ENTRY STRUCTURE", size=8, color=NAVY_900, bold=True)
    p2 = tf.add_paragraph()
    _add_run(p2, "1. MSRA pre-qualification", size=9, color=NAVY_900, bold=True)
    _add_run(p2, " -- yard + workforce + safety credentials", size=9, color=INK)
    p3 = tf.add_paragraph()
    _add_run(p3, "2. MAC-MO IDIQ capture", size=9, color=NAVY_900, bold=True)
    _add_run(p3, " -- multi-award pool, fixed by RMC", size=9, color=INK)
    p4 = tf.add_paragraph()
    _add_run(p4, "3. Fixed-price task orders", size=9, color=NAVY_900, bold=True)
    _add_run(p4, " -- awarded against third-party planner specs", size=9, color=INK)

    add_callout(
        s, R_PANEL_X, Inches(6.05), R_PANEL_W, Inches(0.9),
        "Note:",
        "Entry requires qualification at all three levels AND a physical yard "
        "footprint aligned with one RMC region. The three-tier structure "
        "drives Tier-1 CONUS concentration -- six primes hold ~60% of the "
        "$4.9B segment. Trade-specific IDIQs (HM&E, surface coatings, "
        "insulation, pipe lagging) sum to less than 3% of depot $.")

    add_footer(s,
               "Notes: (1) $4.9B in-scope TAM = $5.0B gross J998/J999 "
               "obligations less $85M FMS carve-out. (2) Contractor tiers: "
               "T1 = CONUS full-ship MSRA holders; T2 = regional yards; "
               "T3 = technical services; T4 = FDNF foreign yards. "
               "(3) IDV scope taxonomy generated via LLM classification of "
               "360 parent IDV descriptions + 471 residual no-IDV awards. "
               "Source: FPDS FY2025 contract obligations on PSCs J998 + J999 "
               "(U.S. Navy + U.S. Coast Guard), classifier-enriched. "
               "Data as of April 2026.")


# ---------- Slide 6: Vessel Mix ----------
def slide_06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Vessel Mix",
              "Surface combatants, amphibious warfare ships, and submarines "
              "accounted for ~62% of FY2025 MRO TAM")

    add_subtitle(s, L_PANEL_X, BODY_TOP, L_PANEL_W,
                 "FY2025 MRO TAM by Vessel Type x Work Segment ($M)")
    add_think_cell_block(
        s, L_PANEL_X, BODY_TOP + Inches(0.32), L_PANEL_W, Inches(4.4),
        "Mekko (Marimekko)",
        ["Work Segment", "Surf Comb.", "Amphib.", "Subs", "Combat Log.",
         "Carriers", "Other"],
        [
            {"cells": ["(column width $M)", "2,120", "1,440", "838", "770",
                      "422", "1,474"], "kind": "sub"},
            {"cells": ["Depot Ship Repair", "1,802", "1,339", "444", "639",
                      "363", "590"]},
            {"cells": ["HM&E", "127", "72", "159", "85", "34", "486"]},
            {"cells": ["Combat Systems", "127", "0", "184", "0", "0", "15"]},
            {"cells": ["Port & Technical", "0", "0", "8", "39", "4", "162"]},
            {"cells": ["Electronics & C4ISR", "85", "14", "42", "0", "21", "206"]},
            {"cells": ["Column total", "2,120", "1,440", "838", "770",
                      "422", "1,474"], "kind": "total"},
        ],
        col_widths=[25, 13, 13, 11, 13, 12, 13],
        right_align_cols={1, 2, 3, 4, 5, 6},
        note="Values shown are computed from the Vessel Type x Work Segment "
             "% stacks; refresh from the Services workbook sheet at final "
             "paste time if precision matters.")

    # Right: two callouts + small vessel-stack table
    add_callout(
        s, R_PANEL_X, BODY_TOP, R_PANEL_W, Inches(1.4),
        "Concentration tracks fleet size and maintenance tempo.",
        "DDG-51 destroyers, LPD / LHD amphibious warships, and Virginia-class "
        "SSNs made up the bulk of FY2025 MRO spend. The three categories "
        "together anchor ~62% of the $7.1B TAM and drive where depot "
        "availabilities cluster geographically (Slide 8).")

    add_callout(
        s, R_PANEL_X, Inches(3.0), R_PANEL_W, Inches(1.55),
        "Submarine spend is structurally understated.",
        "An estimated ~$4-6B of annual nuclear depot work is performed "
        "in-house at the four public naval shipyards (Portsmouth, Norfolk, "
        "Puget Sound, Pearl Harbor) and does not generate FPDS contract "
        "records. Only private-sector OEM and specialty work is captured in "
        "the $838M figure shown here. See Slide 3 for the full scope "
        "reconciliation.")

    # Vessel-category % stack table (reference)
    add_data_table(
        s, R_PANEL_X, Inches(4.75), R_PANEL_W, Inches(2.1),
        ["Category", "Depot", "HM&E", "Combat", "C4ISR", "Port"],
        [
            {"cells": ["Surface Combatants", "85%", "6%", "6%", "4%", "-"]},
            {"cells": ["Amphibious", "93%", "5%", "-", "1%", "-"]},
            {"cells": ["Submarines", "53%", "19%", "22%", "5%", "1%"]},
            {"cells": ["Combat Logistics", "83%", "11%", "-", "-", "5%"]},
            {"cells": ["Aircraft Carriers", "86%", "8%", "-", "5%", "1%"]},
            {"cells": ["Other", "40%", "33%", "1%", "14%", "11%"]},
        ],
        col_widths=[38, 14, 12, 12, 12, 12],
        right_align_cols={1, 2, 3, 4, 5},
        header_size=8, body_size=8, row_height=Inches(0.24))

    add_footer(s,
               "Source: FPDS FY2025 contract obligations (U.S. Navy + U.S. "
               "Coast Guard, 65 services PSCs, post-exclusions). Data as of "
               "April 2026.")


# ---------- Slide 7: Prime Landscape ----------
def slide_07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Prime Landscape",
              "BAE, General Dynamics, and HII collectively captured ~36% of "
              "the FY2025 MRO TAM; the top 10 primes captured ~57%")

    add_subtitle(s, L_PANEL_X, BODY_TOP, L_PANEL_W,
                 "FY2025 MRO Obligations by Contractor - Top 10 ($M) with "
                 "cumulative share (%)")
    add_think_cell_block(
        s, L_PANEL_X, BODY_TOP + Inches(0.32), L_PANEL_W, Inches(4.35),
        "Combination chart (columns + cumulative line)",
        ["Rank", "Contractor", "FY25 $M", "Cum $M", "Cum % of TAM"],
        [
            {"cells": ["1", "BAE", "1,073", "1,073", "15.2%"]},
            {"cells": ["2", "GD", "939", "2,012", "28.5%"]},
            {"cells": ["3", "HII", "516", "2,528", "35.8%"]},
            {"cells": ["4", "Vigor", "440", "2,968", "42.0%"]},
            {"cells": ["5", "Draper", "318", "3,286", "46.5%"]},
            {"cells": ["6", "Detyens", "225", "3,511", "49.7%"]},
            {"cells": ["7", "Epsilon", "149", "3,660", "51.8%"]},
            {"cells": ["8", "East Coast Repair", "142", "3,802", "53.8%"]},
            {"cells": ["9", "Lockheed", "132", "3,934", "55.7%"]},
            {"cells": ["10", "S.C.A.", "112", "4,046", "57.3%"]},
        ],
        col_widths=[8, 32, 18, 20, 22], right_align_cols={2, 3, 4},
        note="Use columns for FY25 $M (left axis 0-1,100) and a line series "
             "for Cum % of TAM (right axis 0-100%). Annotate 'Top 3 = 36%' "
             "and 'Top 10 = 57%' on the line.")

    add_data_table(s, R_PANEL_X, BODY_TOP, R_PANEL_W, Inches(3.1),
                   ["Work Segment", "#1", "#2", "#3"],
                   [
                       {"cells": ["Depot Ship Repair",
                        "BAE  22%", "GD  18%", "HII  10%"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["HM&E",
                        "Global PCCI  7%", "Oceaneering  6%", "HII  6%"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Combat Systems",
                        "Draper  54%", "Lockheed  22%", "Leidos  10%"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Electronics & C4ISR",
                        "Amentum  15%", "L3  15%", "SAIC  15%"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Port & Technical",
                        "S.C.A.  26%", "Waypoint  9%", "Fairlead  6%"],
                        "bolds": [True, False, False, False]},
                   ],
                   col_widths=[34, 22, 22, 22],
                   header_size=9, body_size=8.5, row_height=Inches(0.48))

    add_callout(
        s, R_PANEL_X, Inches(5.2), R_PANEL_W, Inches(1.15),
        "Note:",
        "HII Mission Technologies posted a ~5.0% OI margin in 2025 ($3.0B "
        "revenue / $153M OI). At ~91% service revenue mix, it is the cleanest "
        "pure-services proxy in public comps for the FPDS MRO obligations "
        "captured here.")

    add_footer(s,
               "Notes: (1) BAE includes San Diego, Norfolk, and Jacksonville "
               "ship repair operations. (2) GD includes NASSCO, Electric "
               "Boat, Continental Maritime, and Mission Systems. (3) HII "
               "includes Newport News, Ingalls, Metro Machine, and Marine "
               "Hydraulics International. Cumulative % computed against the "
               "$7,067M Services TAM (Slide 2 denominator). "
               "Sources: FPDS FY2025 contract obligations; SEC 10-K filings "
               "(HII, GD, BWXT), FY2025. Data as of April 2026.")


# ---------- Slide 8: Geographic Context ----------
def slide_08(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Geographic Context",
              "FY2025 depot $ concentrates on the West Coast (SWRMC $1.6B) "
              "and East Coast (MARMC $1.0B); no established RMC anchor sits "
              "on the Gulf Coast")

    add_subtitle(s, L_PANEL_X, BODY_TOP, L_PANEL_W,
                 "FY2025 J998 / J999 depot obligations by RMC region x "
                 "contractor tier ($M)")
    add_think_cell_block(
        s, L_PANEL_X, BODY_TOP + Inches(0.32), L_PANEL_W, Inches(4.35),
        "Mekko (Marimekko)",
        ["Tier", "SWRMC", "MARMC", "Other", "SERMC", "FDNF", "NW RMC", "USCG"],
        [
            {"cells": ["(column width $M)", "1,584", "1,022", "688",
                      "540", "510", "430", "148"], "kind": "sub"},
            {"cells": ["Tier 1 CONUS Complex", "1,140", "680", "420",
                      "350", "0", "330", "95"]},
            {"cells": ["Tier 2 Regional", "310", "200", "150",
                      "120", "0", "60", "35"]},
            {"cells": ["Tier 3 Technical Services", "110", "110", "70",
                      "50", "0", "30", "15"]},
            {"cells": ["Tier 4 FDNF Foreign Yard", "0", "0", "0",
                      "0", "510", "0", "0"]},
            {"cells": ["Other", "24", "32", "48", "20", "0", "10", "3"]},
            {"cells": ["Column total", "1,584", "1,022", "688",
                      "540", "510", "430", "148"], "kind": "total"},
        ],
        col_widths=[26, 11, 11, 10, 11, 10, 10, 11],
        right_align_cols={1, 2, 3, 4, 5, 6, 7},
        note="Mekko column widths are RMC-region $; stacks within each column "
             "are contractor tiers. Tier 4 is 100% of the FDNF column by "
             "construction.")

    add_data_table(s, R_PANEL_X, BODY_TOP, R_PANEL_W, Inches(2.7),
                   ["RMC", "Geography", "Tier-1 Primes in Region",
                    "Navy / USCG Homeports"],
                   [
                       {"cells": ["SWRMC", "San Diego, CA",
                        "BAE SD; GD NASSCO; HII CM + MHI",
                        "San Diego (Pacific Fleet)"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["MARMC", "Norfolk, VA",
                        "BAE Norfolk; HII Metro Machine + MHI",
                        "Norfolk (Atlantic Fleet)"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["SERMC", "Mayport FL + Charleston SC",
                        "Detyens; NASSCO Mayport; BAE JAX",
                        "Mayport; Kings Bay; Charleston"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["NW RMC", "Bremerton/Everett + Portland",
                        "Vigor Seattle + Portland",
                        "Bremerton; Everett"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["FDNF", "Yokosuka + Naples + Bahrain",
                        "Hanwha (KR); Navantia (ES); Sumitomo/Mitsubishi/UniThai",
                        "Forward-deployed squadrons"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["USCG SFLC", "Baltimore MD (HQ)",
                        "None (USCG work distributed across T2/T3)",
                        "USCG homeports nationwide"],
                        "bolds": [True, False, False, False]},
                   ],
                   col_widths=[12, 20, 38, 30],
                   header_size=8, body_size=7, row_height=Inches(0.38))

    # Candidate site box
    cand = add_textbox(s, R_PANEL_X, Inches(4.55), R_PANEL_W, Inches(0.9),
                       fill=SLATE_100, line=RULE, line_width=Pt(0.5))
    tf = cand.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    _add_run(p, "CANDIDATE SITE PROXIMITY", size=8, color=NAVY_900, bold=True)
    p2 = tf.add_paragraph()
    _add_run(p2, "Gulf Coast (Brownsville, TX) -> nearest active RMC is "
             "SERMC (Mayport) at ~1,100 nm.", size=8.5, color=INK)
    p3 = tf.add_paragraph()
    _add_run(p3, "California (site TBD) -> would sit within or adjacent to "
             "SWRMC's $1.6B depot pool (or NW RMC $430M).",
             size=8.5, color=INK)

    add_callout(
        s, R_PANEL_X, Inches(5.6), R_PANEL_W, Inches(1.35),
        "Note:",
        "FY25 depot $ concentrates at locations where the Navy already "
        "homeports major surface-combatant and amphibious-ship squadrons "
        "(San Diego, Norfolk, Mayport, Bremerton/Everett). The candidate Gulf "
        "Coast site does not sit within an established RMC footprint, while "
        "the candidate California site would sit within or adjacent to "
        "SWRMC's $1.6B annual depot pool.")

    add_footer(s,
               "Notes: (1) RMC $ and tier composition drawn from the Depot "
               "Ship Repair RMC x Tier crosstab (J998/J999 only, FY2025 "
               "post-FMS). Non-depot Services MRO $ ($2.3B) is not allocated "
               "to RMCs. (2) FDNF = Forward Deployed Naval Forces; Tier 4 is "
               "100% of FDNF $ by construction. (3) Brownsville, TX to "
               "Mayport, FL is approximately 1,100 nautical miles. "
               "Source: FPDS FY2025 contract obligations on PSCs J998 + J999, "
               "classifier-enriched with RMC geography tags. "
               "Data as of April 2026.")


# ---------- Slide 9: Appropriation Sourcing ----------
def slide_09(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Appropriation Sourcing",
              "The $7.4B FY2025 MRO-PSC universe is funded across ~10 federal "
              "appropriations; OMN and OPN each carry ~35-37%, with "
              "Defense-Wide RDT&E a surprising 11%")

    add_subtitle(s, L_PANEL_X, BODY_TOP, L_PANEL_W,
                 "FY2025 MRO-PSC obligations by Treasury appropriation "
                 "($M, TAS-attributed); OPN drilled to Budget Activity")

    # Two separate think-cell data blocks, side by side
    split_w = (L_PANEL_W - Inches(0.25)) / 2
    add_think_cell_block(
        s, L_PANEL_X, BODY_TOP + Inches(0.55), split_w, Inches(4.2),
        "Stacked Column (Col 1 of 2)",
        ["Appropriation", "FY25 $M", "% of TAS"],
        [
            {"cells": ["OMN", "2,761", "37.3%"], "bolds": [True, False, False]},
            {"cells": ["OPN", "2,588", "35.0%"], "bolds": [True, False, False]},
            {"cells": ["RDT&E, Defense-Wide", "780", "10.5%"]},
            {"cells": ["USCG (OE + AC&I + other)", "320", "4.3%"]},
            {"cells": ["Defense-Wide other", "307", "4.2%"]},
            {"cells": ["Navy other (APN + WPN)", "295", "4.0%"]},
            {"cells": ["Air Force", "165", "2.2%"]},
            {"cells": ["Army", "131", "1.8%"]},
            {"cells": ["SCN + other agency", "53", "0.7%"]},
            {"cells": ["TAS Total", "7,400", "100%"], "kind": "total"},
        ],
        col_widths=[50, 25, 25], right_align_cols={1, 2},
        note="Single stacked column; order segments largest-to-smallest from "
             "bottom.")

    add_think_cell_block(
        s, L_PANEL_X + split_w + Inches(0.25), BODY_TOP + Inches(0.55),
        split_w, Inches(4.2),
        "Stacked Column (Col 2 of 2 - OPN drill)",
        ["OPN Budget Activity", "FY25 $M", "% of OPN"],
        [
            {"cells": ["BA-7 Personnel & Cmd Support Equip", "1,591", "61.5%"],
             "bolds": [True, False, False]},
            {"cells": ["BA-8 Spares & Repair Parts", "825", "31.9%"]},
            {"cells": ["Other BAs (BA-1 + Undistributed)", "171", "6.6%"]},
            {"cells": ["OPN Total", "2,587", "100%"], "kind": "total"},
        ],
        col_widths=[50, 25, 25], right_align_cols={1, 2},
        note="Second chart visually anchors to the OPN segment in the Col 1 "
             "chart via a bracket/arrow drawn after insertion.")

    # Right: four takeaway callouts
    yy = BODY_TOP
    takeaways = [
        ("OMN + OPN = 72% of MRO $",
         "Navy sustainment drives the market. The appropriation story is a "
         "Navy story, not a DoD-wide one."),
        ("OPN splits 62% Command Support Equip / 32% Spares",
         "BA-7 funds installation / modernization electronics, C4ISR, and "
         "combat-system integration. BA-8 is spares and consumables. Depot "
         "availabilities are funded through BA-7, not OMN CE-928."),
        ("RDT&E Defense-Wide = $780M (11%)",
         "Almost entirely Trident II / SSP / SMDC sustainment on J-series "
         "PSCs. Draper MK7 Trident ($318M FY25) is the single-largest MRO "
         "PIID."),
        ("SCN on MRO PSCs = ~$40M",
         "De minimis. Nuclear-platform MRO bundles under PSC 1905 "
         "shipbuilding, not onto MRO PSCs - see Slide 3."),
    ]
    for title, body in takeaways:
        add_callout(s, R_PANEL_X, yy, R_PANEL_W, Inches(1.0), title, body)
        yy += Inches(1.1)

    add_callout(
        s, R_PANEL_X, Inches(6.15), R_PANEL_W, Inches(0.75),
        "Note:",
        "Reconciling FPDS FY25 MRO $ to OMN CE-928 alone leaves a naive $4.7B "
        "gap. The gap disappears once OPN, RDT&E-DW, USCG, and the other "
        "appropriation colors are included.")

    add_footer(s,
               "Notes: (1) 49% of FY25 MRO $ directly TAS-attributed from "
               "Treasury File C (USAspending funding endpoint); 51% imputed "
               "via per-PSC-bucket appropriation ratios. (2) $7,400M TAS "
               "Total is pre-exclusion MRO-PSC universe, ~$333M larger than "
               "the $7,067M Services TAM. "
               "Source: USAspending /awards/funding/ joined to FPDS FY2025 "
               "contract obligations. Data as of April 2026.")


# ---------- Slide 10: TAM Framing ----------
def slide_10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "TAM Framing",
              "The $7.1B FY25 contracting-market TAM apportions to a "
              "~$3.3-3.8B range of FY25-delivered contractor revenue under "
              "period-of-performance sensitivity; the deck presents Frame A "
              "(contracting activity) by industry convention")

    add_subtitle(s, L_PANEL_X, BODY_TOP, L_PANEL_W,
                 "FY2025 Services MRO TAM under alternative "
                 "period-of-performance apportionment methods ($M)")
    add_think_cell_block(
        s, L_PANEL_X, BODY_TOP + Inches(0.55), L_PANEL_W, Inches(3.1),
        "Combination chart (columns + reference line)",
        ["Rank", "Method", "FY25 $M", "% of Frame A"],
        [
            {"cells": ["1", "M5 Frame A - no apportionment", "7,067", "100.0%"],
             "bolds": [False, True, False, False], "kind": "total"},
            {"cells": ["2", "M2 12-month POP cap", "3,747", "53.0%"]},
            {"cells": ["3", "M3 18-month POP cap", "3,499", "49.5%"]},
            {"cells": ["4", "M1 Pure Linear (primary Frame B)", "3,354", "47.5%"]},
            {"cells": ["5", "M4 Front-loaded 60/40", "3,266", "46.2%"]},
        ],
        col_widths=[8, 47, 20, 25], right_align_cols={2, 3},
        note="Columns: FY25 $M (left axis 0-8,000). Add a horizontal "
             "reference line at $3,500M labeled 'Central Frame B ~$3.5B "
             "(49% of Frame A)'.")

    add_data_table(s, R_PANEL_X, BODY_TOP, R_PANEL_W, Inches(2.4),
                   ["Work Segment", "Frame A $M", "Frame B M1 $M", "Apport %"],
                   [
                       {"cells": ["Port & Technical Services", "431", "307", "71.3%"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Depot Ship Repair", "4,781", "2,468", "51.6%"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["HM&E", "938", "363", "38.7%"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Electronics & C4ISR", "333", "105", "31.4%"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Combat Systems Sustainment", "585", "112", "19.1%"],
                        "bolds": [True, False, False, False]},
                       {"cells": ["Total", "7,067", "3,354", "47.5%"],
                        "kind": "total"},
                   ],
                   col_widths=[42, 20, 20, 18],
                   right_align_cols={1, 2, 3},
                   header_size=9, body_size=8.5, row_height=Inches(0.35))

    # Reading-the-table box
    read_box = add_textbox(s, R_PANEL_X, Inches(4.55), R_PANEL_W, Inches(1.3),
                           fill=SLATE_100, line=RULE, line_width=Pt(0.5))
    tf = read_box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    _add_run(p, "READING THE SEGMENT TABLE", size=8, color=NAVY_900, bold=True)
    p2 = tf.add_paragraph()
    _add_run(p2, "- Combat Systems apportions at only 19% -- Trident II / "
             "AEGIS integration runs on 5-10 year POPs (Draper MK7 LE2 alone "
             "is $318M, POP FY24 -> FY29).", size=8.5, color=INK)
    p3 = tf.add_paragraph()
    _add_run(p3, "- Port & Technical apportions at 71% -- husbanding and QC "
             "work is short-cycle and POP aligns with FY25.",
             size=8.5, color=INK)
    p4 = tf.add_paragraph()
    _add_run(p4, "- Depot Ship Repair sits near the blended average at ~52%.",
             size=8.5, color=INK)

    add_callout(
        s, R_PANEL_X, Inches(5.95), R_PANEL_W, Inches(0.95),
        "Note:",
        "Frame A (contracting activity) is the industry convention for "
        "federal contracting TAM in PE / M&A / sell-side diligence. Frame B "
        "(revenue earned) is closer to a contractor's reported P&L but is not "
        "how federal markets are typically sized. Readers modeling a specific "
        "contractor can multiply the relevant segment $ by the apportionment "
        "rate above for a first-order Frame B estimate.")

    add_footer(s,
               "Notes: (1) FY25 window: 2024-10-01 to 2025-09-30. POP "
               "apportionment uses start_date / end_date fields on FPDS award "
               "rows (100% populated in the MRO dataset). (2) M1 Pure Linear "
               "is the primary Frame B estimate; M2-M4 are bounds. (3) True "
               "Frame B would require contractor-side billing / earned-value "
               "data unavailable in federal sources. "
               "Source: FPDS FY2025 contract obligations (U.S. Navy + U.S. "
               "Coast Guard, 65 services PSCs, post-exclusions). Period-of-"
               "performance dates per FPDS. Data as of April 2026.")


# ---------- Main ----------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)

    out = Path(__file__).parent / "MRO_Deck_v0.1.pptx"
    prs.save(out)
    print(f"Saved {out}")


if __name__ == "__main__":
    main()
